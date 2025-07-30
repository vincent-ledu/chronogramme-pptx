# 📦 Imports
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json
import argparse
import re
import hashlib

# 🛠️ Argument parser
parser = argparse.ArgumentParser(description="Génère des slides chronogrammes par tribue.")
parser.add_argument("excel_file", help="Fichier Excel des données")
parser.add_argument("--config", default="config.json", help="Fichier de configuration JSON")
parser.add_argument("--out", default=".", help="Répertoire de sortie pour les fichiers PPTX")
parser.add_argument("--template", default="exemple_chronogramme.pptx", help="Powerpoint modèle pour les slides")
args = parser.parse_args()

# 📂 Chargement des fichiers
excel_path = args.excel_file
config_path = args.config
output_dir = args.out
template_path = args.template

# 📁 Vérification des chemins
if not os.path.exists(template_path):
    print(f"❌ Modèle PowerPoint introuvable : {template_path}")
    sys.exit(1)
if not os.path.exists(excel_path):
    print(f"❌ Fichier Excel introuvable : {excel_path}")
    sys.exit(1)
if not os.path.exists(output_dir):
    print(f"📂 Répertoire de sortie inexistant, création : {output_dir}")
    os.makedirs(output_dir)
if not os.path.exists(config_path):
    print(f"❌ Fichier de configuration introuvable : {config_path}")
    sys.exit(1)

with open(config_path, "r", encoding="utf-8") as f:
    config = json.load(f)

# 📊 Noms des colonnes
col_produit = config["colonne_produit"]
col_solution = config["colonne_solution"]
col_planif = config["colonne_planification"]
col_tribue = config["colonne_tribue"]
col_squad = config["colonne_squad"]
col_kube = config["colonne_full_kube"]
col_z = config["colonne_full_z"]
col_mosart = config["colonne_mosart"]
col_critique = config["colonne_critique"]
col_validate = config["colonne_validate"]
col_decom = config["colonne_decom"]
col_type = config["colonne_type"]
col_realise = config["colonne_realise"]

# 🎨 Icônes utilisées
icone_kube = "icones/" + config["icone_kube"]
icone_critique = "icones/" + config["icone_critique"]
icone_check_green = "icones/" + config["icone_check_green"]
icone_check_blue = "icones/" + config["icone_check_blue"]
icone_z = "icones/" + config["icone_z"]
icone_reconstruction = "icones/" + config["icone_reconstruction"]
icone_restauration = "icones/" + config["icone_restauration"]
icone_resynchro = "icones/" + config["icone_resynchro"]

# 🌈 Conversion couleurs squad
def hex_to_rgbcolor(hex_code):
    hex_code = hex_code.lstrip("#")
    return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))

squad_colors_raw = config.get("squad_colors", {})
squad_color_map = {
    squad: hex_to_rgbcolor(hex_code)
    for squad, hex_code in squad_colors_raw.items()
}

# 🗓️ Positions temporelles des trimestres
positions = {
    "T1/2025": Inches(1.0), "T2/2025": Inches(2.5), "T3/2025": Inches(4.0), "T4/2025": Inches(5.5),
    "T1/2026": Inches(7.0), "T2/2026": Inches(8.5), "T3/2026": Inches(10.0), "T4/2026": Inches(11.5),
}
positions_valides = set(positions.keys())

# 📥 Lecture des données Excel
df = pd.read_excel(excel_path, keep_default_na=False)
tribues = df[col_tribue].dropna().unique()

# 🔢 Fonction pour trier les trimestres
def trimestre_to_sort_key(trimestre):
    match = re.match(r"T([1-4])/(\d{4})", str(trimestre).strip())
    if match:
        t, y = match.groups()
        return int(y) * 10 + int(t)
    return float('inf')

# 🔄 Fonction pour dupliquer les lignes critiques
def dupliquer_lignes_critiques(df):
    lignes_critiques = df[df[col_critique] == 1].copy()
    def trimestre_plus_1(val):
        match = re.match(r"T([1-4])/(\d{4})", str(val).strip())
        if match:
            tri, annee = match.groups()
            return f"T{tri}/{int(annee) + 1}"
        return val
    lignes_critiques[col_planif] = lignes_critiques[col_planif].apply(trimestre_plus_1)
    return pd.concat([df, lignes_critiques], ignore_index=True)

def generate_color_from_string(s):
    h = hashlib.md5(s.encode()).hexdigest()
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    # optionnel : éclaircir un peu
    if r + g + b < 300:
        r, g, b = min(255, r + 60), min(255, g + 60), min(255, b + 60)
    return RGBColor(r, g, b)

# 🌀 Boucle principale par tribue
for tribue in tribues:
    df_tribue = df[df[col_tribue] == tribue].copy()

    # 🧼 Nettoyage des colonnes booléennes
    bool_cols = [col_kube, col_critique, col_mosart, col_decom, col_validate, col_z]
    for col in bool_cols:
        df_tribue[col] = df_tribue[col].astype(str).str.strip().str.lower().map(lambda x: x == "oui" or x.startswith("lot")).astype(int)

    # 🧱 Nettoyage des colonnes texte
    df_tribue[col_type] = df_tribue[col_type].fillna("")
    df_tribue[col_realise] = df_tribue[col_realise].fillna("")

    # 🚫 Exclusion des lignes avec "NA" ou "NR"
    initial_count = len(df_tribue)
    df_tribue = df_tribue[~df_tribue[col_realise].astype(str).str.strip().str.upper().isin(["NR", "NA"])]
    exclues_count = initial_count - len(df_tribue)
    
    if exclues_count > 0:
        print(f"ℹ️  {exclues_count} lignes(s) exclue(s) pour '{tribue}' car 'réalisé' = NA ou NR")

    # 🔢 Tri temporel
    df_tribue["__sort_key"] = df_tribue[col_planif].apply(trimestre_to_sort_key)

    # 📊 Fusion des lignes produit/solution
    fusionnees = (
        df_tribue.sort_values("__sort_key")
        .groupby([col_produit, col_solution])
        .agg({
            col_planif: "first",
            col_tribue: "first",
            col_squad: "first",
            col_kube: "max",
            col_mosart: "max",
            col_critique: "max",
            col_decom: "max",
            col_validate: "max",
            col_z: "max",
            col_type: lambda x: list(x),
            col_realise: lambda x: list(x)
        })
        .reset_index()
    )

    # 🔄 Dupliquer les lignes critiques
    # df_tribue = dupliquer_lignes_critiques(df_tribue)

    # 🖼️ Chargement du modèle PowerPoint
    prs = Presentation("exemple_chronogramme.pptx")
    slide = prs.slides[0]
    slide_invalide = prs.slides[1]

    # 📐 Mise en page
    height = Inches(0.22)
    width = Inches(1.0)
    base_top = Inches(1.5)
    vspace = Inches(0.36)
    ligne_par_trimestre = {}
    lignes_inconnues = 0

    # 🧱 Boucle sur chaque ligne fusionnée
    for _, row in fusionnees.iterrows():
        produit = str(row[col_produit])
        solution = str(row[col_solution])
        trimestre = str(row[col_planif]).strip()
        squad = str(row[col_squad]).strip()
        full_kube = row[col_kube] == 1
        mosart = row[col_mosart] == 1
        decom = row[col_decom] == 1
        critique = row[col_critique] == 1
        validate = row[col_validate] == 1
        full_z = row[col_z] == 1
        types = [str(t).strip().lower() for t in row[col_type]]
        realises = [str(r).strip().lower() for r in row[col_realise]]

        types_realises = {
            t: (realises[i] == "oui" if i < len(realises) else False)
            for i, t in enumerate(types)
        }

        if squad in squad_color_map:
            color = squad_color_map[squad]
        else:
            color = generate_color_from_string(squad)
            squad_color_map[squad] = color  # mémorise pour réutilisation

        complement = ""

        # 🧭 Positionnement sur le slide
        if trimestre in positions_valides:
            left = positions[trimestre]
            ligne = ligne_par_trimestre.get(trimestre, 0)
            top = base_top + vspace * ligne
            ligne_par_trimestre[trimestre] = ligne + 1
            target_slide = slide
        else:
            left = Inches(1.0)
            top = Inches(1.0 + lignes_inconnues * 0.4)
            lignes_inconnues += 1
            target_slide = slide_invalide
            complement = f"\n{trimestre}"
            print(f"⚠️  {produit}-{solution} n'est pas sur un trimestre valide : {trimestre}")

        # 🧱 Boîte principale
        textbox = target_slide.shapes.add_shape(1, left, top, width, height)
        textbox.text = f"{produit}-{solution}{complement}"

        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = color

        line = textbox.line
        line.color.rgb = RGBColor(0, 0, 0)
        if mosart:
            line.width = Pt(2.5)
            line.color.rgb = RGBColor(255, 102, 0)
        if decom:
            line.width = Pt(2.5)
            line.color.rgb = RGBColor(80, 80, 80)

        text_frame = textbox.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(8)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

        # 🎯 Ajout des icônes de type (non réalisés)
        bot_top = top + Inches(0.12)

        if "reconstruction" in types and not types_realises.get("reconstruction", False):
            target_slide.shapes.add_picture(icone_reconstruction, left - Inches(0.0), bot_top, Inches(0.2), Inches(0.2))
        if "restauration bdd" in types and not types_realises.get("restauration bdd", False):
            target_slide.shapes.add_picture(icone_restauration, left + width / 2 - Inches(0.1), bot_top, Inches(0.2), Inches(0.2))
        if "resynchronisation" in types and not types_realises.get("resynchronisation", False):
            target_slide.shapes.add_picture(icone_resynchro, left + width - Inches(0.2), bot_top, Inches(0.2), Inches(0.2))

        # ✅ Validation
        if validate:
            target_slide.shapes.add_picture(icone_check_blue, left + width / 2 - Inches(0.22), bot_top, Inches(0.22), Inches(0.22))
        if all(types_realises.get(t, False) for t in types if t):
            target_slide.shapes.add_picture(icone_check_green, left + width / 2 + Inches(0.22), bot_top, Inches(0.22), Inches(0.22))

        # 🔧 Autres icônes spécifiques
        if full_kube:
            target_slide.shapes.add_picture(icone_kube, left - Inches(0.15), top, Inches(0.22), Inches(0.22))
        if full_z:
            target_slide.shapes.add_picture(icone_z, left - Inches(0.15), top, Inches(0.22), Inches(0.22))
        if critique:
            target_slide.shapes.add_picture(icone_critique, left + width - Inches(0.15), top + Inches(0.02), Inches(0.22), Inches(0.22))

    # 📘 Légende
    legend_left = Inches(0.5)
    legend_top_base = Inches(5.5)
    legend_height = Inches(0.3)
    legend_width = Inches(3.0)
    legend_vspace = Inches(0.32)

    squads_uniques = sorted(fusionnees[col_squad].dropna().unique())
    for i, squad in enumerate(squads_uniques):
        top = legend_top_base + i * legend_vspace
        box = slide.shapes.add_shape(1, legend_left, top, legend_width, legend_height)
        box.text = squad
        box.fill.solid()
        box.fill.fore_color.rgb = squad_color_map.get(squad, RGBColor(200, 200, 200))
        box.line.color.rgb = RGBColor(0, 0, 0)
        text_frame = box.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    # 💾 Sauvegarde du fichier
    date_str = pd.Timestamp.now().strftime("%Y%m%d")
    nom_fichier = f"{output_dir}/{date_str}_chronogramme_{tribue.replace(' ', '_')}.pptx"
    prs.save(nom_fichier)
    print(f"✅ Fichier généré pour {tribue} : {nom_fichier}")